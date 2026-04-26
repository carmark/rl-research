#!/usr/bin/env python3
"""Generate Agentic RL Infra comparison PPT (dark theme, ~25 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x1a, 0x1b, 0x2e)
BG_LIGHT = RGBColor(0x24, 0x25, 0x3d)
ACCENT   = RGBColor(0x4a, 0x9e, 0xff)  # blue
WHITE    = RGBColor(0xff, 0xff, 0xff)
GRAY     = RGBColor(0xaa, 0xaa, 0xcc)
LIGHT_GRAY = RGBColor(0xdd, 0xdd, 0xee)
DARK_ROW = RGBColor(0x20, 0x21, 0x38)
ALT_ROW  = RGBColor(0x28, 0x29, 0x44)

C_ROLLOUT  = RGBColor(0x4a, 0x9e, 0xff)  # blue
C_TRAIN    = RGBColor(0x4e, 0xc9, 0x6e)  # green
C_DATA     = RGBColor(0xff, 0x9f, 0x43)  # orange
C_IO       = RGBColor(0xb0, 0x6a, 0xff)  # purple
C_YELLOW   = RGBColor(0xff, 0xd9, 0x3d)
C_RED      = RGBColor(0xff, 0x6b, 0x6b)
C_CYAN     = RGBColor(0x00, 0xe5, 0xff)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────────────────────

def _set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, size=18, color=WHITE,
              bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def _add_para(tf, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
              space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p


def _title_bar(slide, text, subtitle=None):
    """Top title bar with accent underline."""
    _add_text(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
              text, size=28, color=WHITE, bold=True)
    # accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    if subtitle:
        _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                  subtitle, size=14, color=GRAY)


def _make_table(slide, rows, cols, data, left, top, width, height,
                col_widths=None, header_color=ACCENT, font_size=10):
    """Create a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # style
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Microsoft YaHei"
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1e, 0x3a, 0x5f)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = C_CYAN
                    paragraph.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_ROW if r % 2 == 1 else ALT_ROW
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = LIGHT_GRAY

    return table


def _add_box(slide, left, top, width, height, title, bullets, title_color=ACCENT,
             bg_color=BG_LIGHT, bullet_size=13):
    """Rounded info box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0x3a, 0x3b, 0x55)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    for b in bullets:
        pp = tf.add_paragraph()
        pp.text = f"  {b}"
        pp.font.size = Pt(bullet_size)
        pp.font.color.rgb = LIGHT_GRAY
        pp.font.name = "Microsoft YaHei"
        pp.space_before = Pt(2)
    return shape


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 ─ Cover
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.2),
          "Agentic RL 训练系统基础设施", size=44, color=WHITE, bold=True,
          align=PP_ALIGN.CENTER)
_add_text(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.8),
          "深度调研与横向对比", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# date + source
_add_text(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5),
          "2026-04-26  |  综合 8 大系统  |  参考综述: The Landscape of Agentic RL for LLMs (TMLR 2026)",
          size=14, color=GRAY, align=PP_ALIGN.CENTER)

# 8 system tags
systems = ["verl (字节)", "ROLL (阿里)", "SLIME (清华)", "MILES (RadixArk)",
           "Forge (MiniMax)", "Seer (月之暗面)", "rl-swarm (Gensyn)", "ThunderAgent (Together)"]
colors = [C_ROLLOUT, C_TRAIN, C_DATA, C_IO, C_YELLOW, C_RED, C_CYAN, ACCENT]
tag_w = Inches(1.35)
tag_h = Inches(0.45)
start_x = Inches(0.8)
for i, (sys, clr) in enumerate(zip(systems, colors)):
    row = i // 4
    col = i % 4
    x = start_x + col * Inches(3.1)
    y = Inches(4.5) + row * Inches(0.7)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), tag_h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x22, 0x23, 0x3c)
    sh.line.color.rgb = clr
    sh.line.width = Pt(2)
    tf = sh.text_frame
    tf.paragraphs[0].text = sys
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = clr
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Accent line at bottom
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(6.6), Inches(9.3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 ─ Background: RLHF → Agentic RL
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "背景: 从 RLHF 到 Agentic RL 的范式转变")

# Three era boxes
eras = [
    ("RLHF 时代 (2023-2024)", C_ROLLOUT,
     ["单步 MDP，同步批处理", "TRL / DeepSpeed-Chat", "短序列 (2K-4K)", "人类偏好标注奖励"]),
    ("RL Scaling 时代 (2025)", C_TRAIN,
     ["异构调度 (verl, OpenRLHF)", "长 CoT 推理 (8K-64K)", "671B MoE 模型支持", "vLLM/SGLang 标准化"]),
    ("Agentic RL 时代 (2025-2026)", C_YELLOW,
     ["POMDP · 多轮交互 · 变长轨迹", "环境延迟不可预测", "200K+ Token 上下文", "CM as Action · 环境泄露防护"]),
]

for i, (title, clr, bullets) in enumerate(eras):
    x = Inches(0.5) + i * Inches(4.2)
    _add_box(s, x, Inches(1.3), Inches(3.9), Inches(2.8), title, bullets,
             title_color=clr, bullet_size=12)

# Arrow between them
for i in range(2):
    x = Inches(4.4) + i * Inches(4.2)
    _add_text(s, x, Inches(2.3), Inches(0.5), Inches(0.5), "→", size=36, color=ACCENT,
              align=PP_ALIGN.CENTER)

# Four subsystems
_add_text(s, Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
          "RL 训练的四大子系统", size=20, color=ACCENT, bold=True)
sub_sys = [
    ("Rollout 生成", C_ROLLOUT, "推理引擎 · 长尾处理\n占总时间 63%-87%"),
    ("Training 优化", C_TRAIN, "并行策略 · 算法实现\nPPO→GRPO→DAPO→CISPO"),
    ("Data Processing", C_DATA, "数据流转 · 奖励计算\n零拷贝 · 前缀共享"),
    ("I/O 通信", C_IO, "权重同步 · 调度\nNCCL · Ray · Gossip"),
]
for i, (name, clr, desc) in enumerate(sub_sys):
    x = Inches(0.5) + i * Inches(3.2)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.0), Inches(2.95), Inches(1.6))
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG_LIGHT
    sh.line.color.rgb = clr
    sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(15); p.font.color.rgb = clr; p.font.bold = True; p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    pp = tf.add_paragraph(); pp.text = desc
    pp.font.size = Pt(11); pp.font.color.rgb = LIGHT_GRAY; pp.font.name = "Microsoft YaHei"
    pp.alignment = PP_ALIGN.CENTER; pp.space_before = Pt(6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 ─ 8 Systems Overview Table
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "调研总览: 8 个 Agentic RL 训练系统")

data = [
    ["系统", "组织", "架构模式", "耦合度", "开源", "论文"],
    ["verl", "字节跳动", "单控制器+SPMD", "共置(可解耦)", "✅ GitHub", "EuroSys 2025"],
    ["ROLL", "阿里巴巴", "单控制器", "灵活(3模式)", "✅ GitHub", "arXiv 多篇"],
    ["SLIME", "清华 THUDM", "编排层解耦", "解耦", "✅ GitHub", "arXiv"],
    ["MILES", "RadixArk", "编排层解耦", "解耦", "✅ GitHub", "arXiv"],
    ["Forge", "MiniMax", "三层中间件", "完全解耦", "❌ 博客", "HF Blog"],
    ["Seer", "月之暗面", "集中式", "动态切换", "❌ 论文", "arXiv"],
    ["rl-swarm", "Gensyn", "P2P 去中心化", "完全解耦", "✅ GitHub", "arXiv"],
    ["ThunderAgent", "Together AI", "Program感知", "工作流解耦", "✅ GitHub", "arXiv 2602"],
]
cw = [Inches(1.6), Inches(1.3), Inches(2.0), Inches(1.8), Inches(1.3), Inches(1.6)]
_make_table(s, 9, 6, data, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5),
            col_widths=cw, font_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 ─ LLM RL vs Agentic RL
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "LLM RL vs. Agentic RL: 核心范式区别",
           "参考: The Landscape of Agentic RL for LLMs (TMLR 2026)")

data = [
    ["维度", "LLM RL (传统)", "Agentic RL"],
    ["MDP 模型", "退化的单步 MDP", "时序扩展 POMDP"],
    ["状态", "Prompt + 已生成 Token", "完整环境状态 (部分可观)"],
    ["动作空间", "词汇表 Token", "Token + 工具 + API + GUI + 代码"],
    ["转移动态", "确定性 (下一 Token 拼接)", "环境介导、随机、可能对抗"],
    ["奖励函数", "偏好模型 (静态)", "任务指标、执行验证、Critique"],
    ["交互", "单轮", "多轮 · 多步 · 持续环境交互"],
    ["信用分配", "Token 级", "跨步骤/跨轮次长程信用分配"],
]
cw = [Inches(1.8), Inches(4.0), Inches(5.5)]
_make_table(s, 8, 3, data, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5),
            col_widths=cw, font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 ─ RL Granularity
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "RL 粒度: Token / Turn / Trajectory / Bi-Level")

data = [
    ["粒度", "描述", "优势", "劣势", "代表系统"],
    ["Token 级", "每个Token=独立MDP步骤", "细粒度信用分配", "无法捕获Agent行为结构", "verl(PPO), SLIME"],
    ["Turn 级", "每个Agent轮次=一个MDP步骤", "匹配任务结构", "错过轮次内部动态", "Forge (Agent多轮)"],
    ["Trajectory 级", "完整轨迹获得单一奖励", "奖励设计简单", "信号稀疏·信用分配难", "rl-swarm (SAPO)"],
    ["Chunk 级", "环境交互间连续段=Chunk", "匹配Agent交互结构", "实现复杂度高", "ROLL (Chunked MDP)"],
    ["双层 Bi-Level", "组合 Turn 级和 Token 级", "结构对齐+细粒度", "复杂度最高", "前沿研究方向"],
]
cw = [Inches(1.5), Inches(3.0), Inches(2.2), Inches(2.5), Inches(2.5)]
_make_table(s, 6, 5, data, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5),
            col_widths=cw, font_size=12)

_add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.5),
          "关键发现: Token 级 MDP 对 Agentic 任务日益不适用。Turn/Step 应成为 LLM Agent 的适当动作表示。",
          size=13, color=C_YELLOW, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 ─ Three Core Challenges
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "Agentic RL Infra 三大核心挑战")

challenges = [
    ("挑战一: Off-Policy 与异步训练", C_RED,
     ["异步训练不可避免引入 Off-Policy 偏差",
      "仅用最新数据 → 采样偏差 (偏向快且简单样本)",
      "仅用最旧数据 → 长尾效应拖累吞吐",
      "解法: 滑动窗口(Forge)·双边IS(SLIME)",
      "    Chunked MDP IS(ROLL)·同步消泡(Seer)"]),
    ("挑战二: 环境管理与数据质量", C_DATA,
     ["环境泄露: Agent中间产物可能提供hint导致作弊",
      "伪阳性: 不完整测试数据导致虚假正向奖励",
      "No-op 通过: 不执行操作即可通过测试",
      "解法: Rock沙箱+环境清理(ROLL)",
      "    LLM-as-Judge验证(ROLL)·Gateway隔离(Forge)"]),
    ("挑战三: 上下文管理", C_IO,
     ["多轮交互产生超长上下文 (200K+ Token)",
      "KV Cache 管理与复用成为瓶颈",
      "KV-Cache Thrashing (7.1x延迟膨胀)",
      "解法: CM as Action(Forge)·Mooncake两层缓存(Seer)",
      "    Program感知防KV Thrashing(ThunderAgent)"]),
]
for i, (title, clr, bullets) in enumerate(challenges):
    x = Inches(0.4) + i * Inches(4.2)
    _add_box(s, x, Inches(1.2), Inches(4.0), Inches(5.5), title, bullets,
             title_color=clr, bullet_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 7-14 ─ Per-System (8 slides)
# ═══════════════════════════════════════════════════════════════════════════════

system_pages = [
    {
        "name": "verl (字节跳动)", "color": C_ROLLOUT,
        "subtitle": "HybridFlow 混合控制器  |  EuroSys 2025  |  v0.7.1",
        "arch": [
            "单控制器 (MPMD) + 多控制器 (SPMD) 统一",
            "ResourcePool: Ray placement group 管理 GPU 集合",
            "create_colocated_worker_cls: Actor/Rollout/Ref 同GPU共置",
            "3D-HybridEngine: 训练↔推理零冗余权重切换 (PP→micro-DP)",
        ],
        "innovations": [
            "DataProto/TransferQueue: 控制流与数据流解耦",
            "Server 模式 + AgentLoop: 支持多轮 Agentic 任务",
            "5D 并行: DP+TP+PP+CP+EP",
            "NIXL: GPU-Direct RDMA P2P 通信",
        ],
        "perf": [
            ("vs DeepSpeed-Chat", "3.67x+"),
            ("异步训练加速", "2.35-2.67x"),
            ("最大模型", "671B MoE"),
            ("权重同步延迟", "<300ms"),
        ],
    },
    {
        "name": "ROLL (阿里巴巴)", "color": C_TRAIN,
        "subtitle": "灵活设备映射 + ROLLMUX/ROLLART 资源优化",
        "arch": [
            "AutoDeviceMapping: 共置/解耦/Serverless 三模式",
            "Rollout Scheduler: 样本生命周期管理",
            "ROLLMUX: 阶段级GPU复用，控制面/数据面解耦",
            "ROLLART: 无状态组件 Serverless 弹性伸缩",
        ],
        "innovations": [
            "Chunked MDP: Chunk级奖励和IS计算 (匹配Agent交互)",
            "AgentServer + Rock Sandbox: 环境管理+防泄露",
            "LLM-as-Judge + No-op 验证: 数据质量保障",
            "Dense Reward: 过程奖励 + 时间奖励 + R2G",
        ],
        "perf": [
            ("ROLLMUX 成本效率", "1.84x"),
            ("RollPacker 加速", "2.03-2.56x"),
            ("ROLLART 加速", "1.35-2.05x"),
            ("权重同步 vs verl", "7.87-8.33x"),
        ],
    },
    {
        "name": "SLIME (清华 THUDM)", "color": C_DATA,
        "subtitle": "SGLang 原生设计  |  GLM-4.5/GLM-5 训练框架  |  v0.2.3",
        "arch": [
            "SGLang 原生集成 + Megatron (mbridge) 训练",
            "SlimeRouter (sgl-router): 多实例负载均衡",
            "StringRadixTrie: Token级前缀缓存",
            "Ray 编排: PlacementGroup + RayTrainGroup",
        ],
        "innovations": [
            "APRIL: 主动部分Rollout (过量供应→早停→回收)",
            "Agent-RL 解耦: 用户无需更改Agent环境",
            "双边 IS 采样修正: Off-Policy 偏差控制",
            "FP8 + DeepEP: 355B MoE 6-7x 推理加速",
        ],
        "perf": [
            ("最大模型", "355B MoE (64×H100)"),
            ("APRIL 端到端加速", "40%"),
            ("FP8+DeepEP 推理", "6-7x"),
            ("多轮 vs 单轮", "一致更优"),
        ],
    },
    {
        "name": "MILES (RadixArk)", "color": C_IO,
        "subtitle": "SLIME 企业级分支  |  MoE 路由一致性解决方案",
        "arch": [
            "SGLang FP8 推理 + Megatron FP8 训练",
            "R3 路由回放桥: 连接推理与训练的路由决策",
            "CUDA IPC 零拷贝: 高效权重同步",
            "在线 SFT MTP: 投机解码草稿层",
        ],
        "innovations": [
            "R3 (Rollout Routing Replay): 解决MoE路由不一致",
            "端到端 FP8: 首个Rollout+Training均FP8的框架",
            "INT4 QAT: 1TB模型单机H200可推理",
            "MrlX: 异步共演化多智能体RL框架",
        ],
        "perf": [
            ("投机解码加速", "25%+"),
            ("权重同步降低", "50%"),
            ("单机模型规模", "1TB (INT4)"),
            ("路由不一致", "94%Token→大幅降低"),
        ],
    },
    {
        "name": "Forge (MiniMax)", "color": C_YELLOW,
        "subtitle": "三层中间件架构  |  M2.5 (230B/10B active)  |  原生 Agentic",
        "arch": [
            "Layer 1 Agent Side: White/Black-Box + Context Manager",
            "Layer 2 Middleware: Gateway + Data Pool + Windowed FIFO",
            "Layer 3 Engine: LLM Rollout + Train + L3 KV Cache Pool",
            "PD 解耦: 异构 Prefill-Decode 分离",
        ],
        "innovations": [
            "CISPO: 全Token梯度参与 (vs PPO部分Token被过滤)",
            "Prefix Tree Merging: ~40x 训练前向加速",
            "CM as Agent Action: 上下文管理建模为动作",
            "滑动窗口 (Windowed FIFO): Off-Policy平衡",
        ],
        "perf": [
            ("CISPO vs DAPO", "2x 收敛"),
            ("Prefix Tree 加速", "~40x"),
            ("模型规模", "230B (10B active)"),
            ("上下文支持", "200K Token"),
        ],
    },
    {
        "name": "Seer (月之暗面 Moonshot AI)", "color": C_RED,
        "subtitle": "严格同步 On-Policy  |  不改算法，仅消除气泡",
        "arch": [
            "集中式控制器 + Megatron 训练",
            "Divided Rollout: GRPO组拆分为8K chunk增量调度",
            "DGDS: 分布式分组草稿服务器 (无需独立草稿模型)",
            "Mooncake KVCache: 两层 DRAM+SSD 分布式缓存",
        ],
        "innovations": [
            "Divided Rollout: 细粒度分解+早期打包+晚期调整",
            "DGDS: 压缩后缀树+组内投机解码 (+30-44%)",
            "任务时间奖励 + Reward-to-Go: Dense Reward",
            "Context-Aware 调度: 近似最优最长优先",
        ],
        "perf": [
            ("Rollout 吞吐量", "+74-97%"),
            ("长尾延迟降低", "-75-93%"),
            ("DGDS 贡献", "+30-44%"),
            ("Divided Rollout", "+27-35%"),
        ],
    },
    {
        "name": "rl-swarm (Gensyn)", "color": C_CYAN,
        "subtitle": "完全去中心化 P2P RL 训练  |  消费级 GPU 民主化",
        "arch": [
            "无中心协调: P2P Gossip 网络",
            "Hivemind 协议: 交换Rollout/反馈/批评",
            "NoLoCo: 低通信 Gossip 替代 AllReduce",
            "Gensyn Testnet 智能合约链上协调",
        ],
        "innovations": [
            "SAPO: 共享Rollout文本(非梯度)，架构无关",
            "P2P 训练: 消费级GPU (RTX 3090+) 可参与",
            "'Aha moments' 通过 Swarm 传播",
            "不同架构和规模的模型可共享经验",
        ],
        "perf": [
            ("累计奖励提升", "94% (8×0.5B)"),
            ("社区参与", "数千成员"),
            ("最低硬件", "RTX 3090"),
            ("局限", "大模型收益递减"),
        ],
    },
    {
        "name": "ThunderAgent (Together AI)", "color": ACCENT,
        "subtitle": "Program-Aware Agentic 推理系统  |  arXiv 2602.13692",
        "arch": [
            "Program 抽象: 工作流级持久调度单元",
            "元数据跟踪: 工作流ID/执行阶段/调度状态/Token数",
            "调度与执行后端 (vLLM/SGLang) 完全解耦",
            "可与 Slime/SkyRL 等 RL 框架集成",
        ],
        "innovations": [
            "KV-Cache 感知调度: 防止活跃工作流KV被驱逐",
            "Program 抽象: 从请求级提升到工作流级调度",
            "异步环境准备: 工具调用期间预准备推理环境",
            "解决 KV-Cache Thrashing (7.1x延迟膨胀→近零)",
        ],
        "perf": [
            ("RL Rollout 吞吐量", "1.8-3.9x"),
            ("推理服务吞吐量", "1.5-3.6x"),
            ("KV Thrashing 优化", "7.1x→近零"),
            ("磁盘内存节省", "最高 4.2x"),
        ],
    },
]

for sp in system_pages:
    s = new_slide()
    _title_bar(s, sp["name"], sp["subtitle"])

    # Left: Architecture
    _add_box(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(2.6),
             "架构核心", sp["arch"], title_color=sp["color"], bullet_size=12)

    # Right: Innovations
    _add_box(s, Inches(6.8), Inches(1.2), Inches(6.1), Inches(2.6),
             "关键创新", sp["innovations"], title_color=C_YELLOW, bullet_size=12)

    # Bottom: Performance
    perf_data = [["指标", "数值"]] + [[k, v] for k, v in sp["perf"]]
    _make_table(s, len(perf_data), 2, perf_data,
                Inches(0.4), Inches(4.2), Inches(5.5), Inches(2.8),
                col_widths=[Inches(3.0), Inches(2.5)], font_size=12)

    # Bottom right: key highlight
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(6.8), Inches(4.2), Inches(6.1), Inches(2.8))
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG_LIGHT
    sh.line.color.rgb = sp["color"]
    sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "性能亮点"
    p.font.size = Pt(16); p.font.color.rgb = sp["color"]; p.font.bold = True
    p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
    for k, v in sp["perf"]:
        pp = tf.add_paragraph()
        pp.text = f"{k}: {v}"
        pp.font.size = Pt(14); pp.font.color.rgb = LIGHT_GRAY
        pp.font.name = "Microsoft YaHei"; pp.alignment = PP_ALIGN.CENTER
        pp.space_before = Pt(8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 ─ Rollout Comparison
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "横向对比: Rollout / Generation", "Rollout 占 RL 训练总时间 63%-87%")

data = [
    ["系统", "推理引擎", "推理模式", "投机解码", "长尾处理", "关键优化"],
    ["verl", "vLLM/SGLang\nTRT-LLM", "Server模式", "-", "AgentLoop\n多轮", "动态批处理\n2.35-2.67x"],
    ["ROLL", "vLLM/SGLang", "动态FP8", "-", "Scheduler\n+RollPacker", "RollPacker\n2.03-2.56x"],
    ["SLIME", "SGLang", "HTTP Server", "-", "APRIL\n部分Rollout", "FP8+DeepEP\n6-7x"],
    ["MILES", "SGLang(FP8)", "HTTP Server", "在线SFT\nMTP(25%+)", "APRIL\n+过采样", "R3路由\n记录"],
    ["Forge", "自研", "PD解耦", "Dynamic\nMTP", "Windowed\nFIFO", "L3 KV Cache\n~40x Prefix"],
    ["Seer", "vLLM", "Divided\nRollout", "DGDS\n(30-44%)", "分段消除\n气泡", "+74-97%\n吞吐量"],
    ["rl-swarm", "vLLM/本地", "本地推理", "-", "-", "P2P共享"],
    ["ThunderAgent", "vLLM/SGLang", "Program\n-Aware", "-", "KV-Cache\n感知", "防KV\nThrashing"],
]
cw = [Inches(1.5), Inches(1.6), Inches(1.5), Inches(1.5), Inches(1.6), Inches(1.7)]
_make_table(s, 9, 6, data, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8),
            col_widths=cw, font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 ─ Training Comparison
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "横向对比: Training / Optimization")

data = [
    ["系统", "训练引擎", "并行策略", "核心算法", "训练模式", "量化"],
    ["verl", "FSDP/Megatron", "5D并行\nDP+TP+PP+CP+EP", "PPO/GRPO\nDAPO/DPO", "同步+异步", "实验FP8"],
    ["ROLL", "Megatron/\nFSDP/DS", "自动选择", "PPO/GRPO\nDAPO", "同步+异步\n(Flash)", "FP8\nRollout"],
    ["SLIME", "Megatron\n(mbridge)", "TP+PP+DP\n+CP+EP", "PPO/GRPO\nDAPO", "同步+异步\n(双边IS)", "-"],
    ["MILES", "Megatron\n(mbridge)", "TP+PP+DP+EP", "GRPO+\nTIS/MIS", "同步+异步", "FP8 E2E\nINT4 QAT"],
    ["Forge", "自研(Magi)", "自研并行", "CISPO", "混合域统一", "-"],
    ["Seer", "Megatron", "DP+TP", "PPO/GRPO", "严格同步\nOn-Policy", "-"],
    ["rl-swarm", "trl+\nHivemind", "P2P\nAllReduce", "SAPO\n(GRPO改进)", "去中心化\n异步", "-"],
    ["ThunderAgent", "依赖集成\n框架", "依赖集成\n框架", "依赖集成\n框架", "依赖集成\n框架", "-"],
]
cw = [Inches(1.5), Inches(1.6), Inches(2.0), Inches(1.6), Inches(1.8), Inches(1.3)]
_make_table(s, 9, 6, data, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8),
            col_widths=cw, font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 ─ I/O Comparison
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "横向对比: I/O 与通信")

data = [
    ["系统", "调度框架", "集合通信", "推理通信", "权重同步方案", "同步性能"],
    ["verl", "Ray", "NCCL", "gRPC/HTTP", "3D-HybridEngine\n/NIXL", "<300ms"],
    ["ROLL", "Ray", "NCCL", "gRPC/HTTP", "ROLLMUX\n挂起/恢复", "7.87-8.33x\nvs verl"],
    ["SLIME", "Ray", "NCCL", "HTTP\n(SGLang)", "CUDA IPC\n/分布式", "BF16~48s\nFP8~100s"],
    ["MILES", "Ray", "NCCL", "HTTP\n(SGLang)", "CUDA IPC\n零拷贝", "50%↓"],
    ["Forge", "自研\n中间件", "自研", "OpenAI\nAPI", "L3 KV Cache\nPool", "-"],
    ["Seer", "自研", "NCCL", "内部RPC", "Mooncake\nCheckpoint", "-"],
    ["rl-swarm", "Hivemind", "P2P\nAllReduce", "libp2p", "Gossip\n广播", "高(互联网)"],
    ["ThunderAgent", "Program\n-Aware", "依赖后端", "依赖后端", "-", "1.8-3.9x\nRollout"],
]
cw = [Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.0), Inches(1.5)]
_make_table(s, 9, 6, data, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8),
            col_widths=cw, font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 ─ Agentic Challenges Comparison
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "横向对比: Agentic 特有挑战")

data = [
    ["挑战", "verl", "ROLL", "SLIME", "MILES", "Forge", "Seer", "rl-swarm", "Thunder\nAgent"],
    ["Off-Policy\n控制", "异步+\n版本管理", "Chunked\nMDP IS", "双边IS\n修正", "TIS/MIS", "滑动窗口", "严格同步", "Gossip\n异步", "依赖\n集成"],
    ["环境泄露\n防护", "-", "✅严格\n清理+隔离", "-", "-", "Gateway\n隔离", "-", "-", "Program\n隔离"],
    ["上下文\n管理", "动态\n批处理", "Scheduler", "RadixTrie", "RadixTrie", "CM as\nAction", "Mooncake\n分布式", "本地", "防KV\nThrashing"],
    ["Dense\nReward", "-", "✅过程+\n时间+R2G", "-", "-", "Critique", "✅时间\n+R2G", "-", "-"],
    ["KV Cache", "vLLM/\nSGLang", "vLLM/\nSGLang", "SGLang\nRadix", "SGLang\nRadix", "L3全局\nPool(DFS)", "Mooncake\n(DRAM+SSD)", "本地", "防驱逐\n(7.1x优化)"],
]
cw = [Inches(1.3)] + [Inches(1.3)] * 8
_make_table(s, 7, 9, data, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8),
            col_widths=cw, font_size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 ─ Innovation Quick Reference
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "创新点速查表")

data = [
    ["系统", "核心创新", "解决的问题"],
    ["verl", "HybridFlow + 3D-HybridEngine + TransferQueue", "灵活性与性能统一·控制流数据流解耦"],
    ["ROLL", "ROLLMUX阶段复用 + ROLLART Serverless\n+ Chunked MDP + AgentServer", "GPU利用率·Agentic奖励建模·环境管理"],
    ["SLIME", "APRIL主动部分Rollout + SGLang原生\n+ Agent-RL解耦 + 双边IS", "长尾延迟·Agent生态兼容·Off-Policy"],
    ["MILES", "R3路由回放 + 端到端FP8\n+ 在线SFT投机解码 + MrlX", "MoE路由一致性·量化加速·多智能体"],
    ["Forge", "CISPO + Prefix Tree Merging\n+ CM as Action + 滑动窗口", "全Token梯度·40x加速·上下文感知"],
    ["Seer", "Divided Rollout + DGDS\n+ Mooncake KVCache + 时间奖励", "长CoT气泡消除·Dense Reward"],
    ["rl-swarm", "SAPO P2P训练 + Gossip共享", "去中心化民主化训练"],
    ["ThunderAgent", "Program抽象 + KV-Cache感知调度\n+ 异步环境准备", "KV-Cache Thrashing消除"],
]
cw = [Inches(1.6), Inches(5.5), Inches(4.5)]
_make_table(s, 9, 3, data, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8),
            col_widths=cw, font_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 ─ Performance Data
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "性能关键数据对比")

data = [
    ["系统", "核心指标", "数值", "条件"],
    ["verl", "vs DeepSpeed-Chat", "3.67x+", "原始基准"],
    ["verl", "异步训练加速", "2.35-2.67x", "128 GPU·Qwen2.5-7B"],
    ["ROLL", "ROLLMUX 成本效率", "1.84x", "vs 标准解耦"],
    ["ROLL", "RollPacker 加速", "2.03-2.56x", "vs verl"],
    ["SLIME", "APRIL 端到端加速", "40%", "吞吐量提升"],
    ["SLIME", "FP8+DeepEP 推理", "6-7x", "GLM4.5-355B"],
    ["MILES", "投机解码加速", "25%+", "在线SFT MTP"],
    ["Forge", "Prefix Tree 加速", "~40x", "训练前向传播"],
    ["Forge", "CISPO vs DAPO", "2x", "Qwen2.5-32B"],
    ["Seer", "Rollout 吞吐量", "+74-97%", "32×8 H800"],
    ["Seer", "长尾延迟降低", "-75-93%", "生产负载"],
    ["ThunderAgent", "RL Rollout 吞吐量", "1.8-3.9x", "vs vLLM+SGLang"],
    ["ThunderAgent", "KV Thrashing 优化", "7.1x→近零", "延迟膨胀消除"],
    ["rl-swarm", "累计奖励提升", "94%", "8×Qwen2.5-0.5B"],
]
cw = [Inches(1.6), Inches(3.0), Inches(1.8), Inches(3.0)]
_make_table(s, 15, 4, data, Inches(0.8), Inches(1.1), Inches(11.7), Inches(6.1),
            col_widths=cw, font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 ─ Architecture Spectrum + Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "架构模式光谱 + 技术栈依赖图")

# Spectrum visualization
_add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.4),
          "集中式 ←━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━→ 去中心化",
          size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

spectrum = [
    ("verl", "单控制器\n共置", C_ROLLOUT),
    ("Seer", "集中式\n动态切换", C_RED),
    ("ROLL", "灵活\n3模式", C_TRAIN),
    ("SLIME", "编排层\n解耦", C_DATA),
    ("MILES", "编排层\n模块化", C_IO),
    ("Forge", "中间件\n完全解耦", C_YELLOW),
    ("Thunder\nAgent", "Program感知\n工作流级", ACCENT),
    ("rl-swarm", "P2P\n无中心", C_CYAN),
]
for i, (name, desc, clr) in enumerate(spectrum):
    x = Inches(0.4) + i * Inches(1.6)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(1.45), Inches(1.2))
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_LIGHT
    sh.line.color.rgb = clr; sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(11); p.font.color.rgb = clr; p.font.bold = True
    p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
    pp = tf.add_paragraph(); pp.text = desc
    pp.font.size = Pt(9); pp.font.color.rgb = LIGHT_GRAY
    pp.font.name = "Microsoft YaHei"; pp.alignment = PP_ALIGN.CENTER

# Tech stack dependency
_add_text(s, Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
          "技术栈依赖图", size=18, color=ACCENT, bold=True)

deps = [
    ("推理引擎", C_ROLLOUT, [
        ("vLLM", "verl, ROLL, Seer, rl-swarm, ThunderAgent"),
        ("SGLang", "verl, ROLL, SLIME, MILES, ThunderAgent"),
        ("自研", "Forge, Seer"),
    ]),
    ("训练引擎", C_TRAIN, [
        ("Megatron-LM", "verl, ROLL, SLIME, MILES, Seer"),
        ("FSDP", "verl, ROLL, MILES"),
        ("自研/其他", "Forge(Magi), rl-swarm(trl)"),
    ]),
    ("调度框架", C_DATA, [
        ("Ray", "verl, ROLL, SLIME, MILES"),
        ("自研", "Forge, Seer"),
        ("其他", "Hivemind(rl-swarm), Program(Thunder)"),
    ]),
    ("集合通信", C_IO, [
        ("NCCL", "verl, ROLL, SLIME, MILES, Seer"),
        ("P2P/Gossip", "rl-swarm"),
        ("自研", "Forge"),
    ]),
]
for i, (cat, clr, items) in enumerate(deps):
    x = Inches(0.3) + i * Inches(3.25)
    y_start = Inches(3.6)
    # Category header
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, Inches(3.1), Inches(0.45))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr
    sh.line.fill.background()
    tf = sh.text_frame
    p = tf.paragraphs[0]; p.text = cat
    p.font.size = Pt(13); p.font.color.rgb = BG; p.font.bold = True
    p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER

    for j, (tech, users) in enumerate(items):
        y = y_start + Inches(0.55) + j * Inches(1.05)
        sh2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.1), Inches(0.95))
        sh2.fill.solid(); sh2.fill.fore_color.rgb = BG_LIGHT
        sh2.line.color.rgb = RGBColor(0x3a, 0x3b, 0x55); sh2.line.width = Pt(1)
        tf2 = sh2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.text = tech
        p2.font.size = Pt(12); p2.font.color.rgb = clr; p2.font.bold = True
        p2.font.name = "Microsoft YaHei"
        pp2 = tf2.add_paragraph(); pp2.text = f"← {users}"
        pp2.font.size = Pt(9); pp2.font.color.rgb = GRAY
        pp2.font.name = "Microsoft YaHei"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 ─ Timeline
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "Infra 演进时间线")

# Timeline bar
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(11.7), Pt(6))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Three periods
periods = [
    ("2024\nRLHF 时代", Inches(1.0), C_ROLLOUT,
     ["TRL/DeepSpeed-Chat", "短序列·同步批处理", "单步MDP·人类偏好", "Token级奖励"]),
    ("2025\nRL Scaling 时代", Inches(5.0), C_TRAIN,
     ["verl/OpenRLHF", "异构调度·长CoT", "671B MoE·GRPO族", "vLLM/SGLang标准化"]),
    ("2025-2026\nAgentic RL 时代", Inches(9.0), C_YELLOW,
     ["Forge异步Data Pool", "SLIME双边IS修正", "ROLL Chunked MDP+AgentServer",
      "Seer Divided Rollout", "ThunderAgent Program抽象",
      "POMDP·多轮·变长轨迹", "CM as Action·环境泄露防护"]),
]

for title, x_pos, clr, items in periods:
    # Circle marker
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x_pos, Inches(2.35), Inches(0.35), Inches(0.35))
    dot.fill.solid(); dot.fill.fore_color.rgb = clr; dot.line.fill.background()

    # Title above
    _add_text(s, x_pos - Inches(0.5), Inches(1.3), Inches(3.5), Inches(1.0),
              title, size=14, color=clr, bold=True, align=PP_ALIGN.CENTER)

    # Items below
    box_h = Inches(0.35) * len(items) + Inches(0.5)
    _add_box(s, x_pos - Inches(0.5), Inches(3.0), Inches(3.5), min(box_h, Inches(4.0)),
             "", items, title_color=clr, bg_color=BG_LIGHT, bullet_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 ─ Key Findings
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "核心发现: 四大依赖总结")

deps_data = [
    ("对 Rollout 的依赖", C_ROLLOUT, [
        "所有系统严重依赖高吞吐推理引擎",
        "vLLM/SGLang 成为事实标准 (6/8系统)",
        "Rollout 占总时间 63%-87%，是核心瓶颈",
        "Agentic 场景: 工具调用 head-of-line blocking",
        "趋势: 异步逐对话Server→工作流级调度",
    ]),
    ("对 Training 的依赖", C_TRAIN, [
        "Megatron-LM 主导 (6/8系统)",
        "FSDP 适合中小规模研究场景",
        "算法: PPO→GRPO(50%资源↓)→DAPO→CISPO",
        "Agentic 新方向: Chunked MDP 匹配交互结构",
        "异步训练 + Off-Policy 控制成新热点",
    ]),
    ("对 Data Processing 的依赖", C_DATA, [
        "零拷贝传输: CUDA IPC/RDMA/NIXL",
        "前缀共享: RadixTrie/Prefix Tree (~40x)",
        "数据质量: LLM-as-Judge / No-op 过滤",
        "奖励演进: 人类标注→自动验证(RLVR)",
        "Dense Reward 在 Agentic 场景愈发重要",
    ]),
    ("对 I/O 的依赖", C_IO, [
        "NCCL 仍是集合通信基础 (6/8系统)",
        "权重同步方案差异最大 (各系统各异)",
        "3D-HybridEngine(verl) vs CUDA IPC(SLIME)",
        "ROLLMUX挂起/恢复(ROLL) vs Mooncake(Seer)",
        "Gossip广播(rl-swarm): 去中心化新路径",
    ]),
]

for i, (title, clr, bullets) in enumerate(deps_data):
    row = i // 2
    col = i % 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.2) + row * Inches(3.0)
    _add_box(s, x, y, Inches(6.2), Inches(2.8), title, bullets,
             title_color=clr, bullet_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 ─ Future Outlook
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "未来展望: 8 大方向")

directions = [
    ("1. Token→Chunk/Turn级MDP", "Chunked MDP(ROLL)·Turn级(Forge)\nBi-Level GAE"),
    ("2. 更长上下文 RL (200K+)", "Mooncake(Seer)·L3 Pool(Forge)\n防KV Thrashing(ThunderAgent)"),
    ("3. 异构硬件支持", "AMD GPU(ROLL/SLIME)\n消费级GPU(rl-swarm)"),
    ("4. 多智能体 RL", "MrlX(MILES)·OpenClaw-RL(SLIME)\nSwarm(rl-swarm)"),
    ("5. 去中心化训练", "SAPO Rollout共享(非梯度)\nrl-swarm 可行性已验证"),
    ("6. 环境可扩展性", "AgentServer+Rock(ROLL)\nAWM/AutoForge 环境合成"),
    ("7. 安全与可信 RL", "环境泄露防护(ROLL)\n分布偏移·工具误用防范"),
    ("8. 推理-训练协同调度", "Program抽象(ThunderAgent)\n工作流级协同调度"),
]

clrs = [C_ROLLOUT, C_TRAIN, C_DATA, C_IO, C_YELLOW, C_RED, C_CYAN, ACCENT]
for i, ((title, desc), clr) in enumerate(zip(directions, clrs)):
    row = i // 4
    col = i % 4
    x = Inches(0.3) + col * Inches(3.25)
    y = Inches(1.2) + row * Inches(3.0)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.1), Inches(2.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_LIGHT
    sh.line.color.rgb = clr; sh.line.width = Pt(2)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.color.rgb = clr; p.font.bold = True
    p.font.name = "Microsoft YaHei"
    pp = tf.add_paragraph(); pp.text = desc
    pp.font.size = Pt(11); pp.font.color.rgb = LIGHT_GRAY
    pp.font.name = "Microsoft YaHei"; pp.space_before = Pt(10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 ─ Appendix: RL Algorithm Family Tree + References
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
_title_bar(s, "附录: RL 算法族谱 + 参考文献")

# Algorithm family tree table
data = [
    ["算法族", "特点", "系统覆盖"],
    ["PPO (需Critic)", "裁剪策略比率, KL惩罚", "verl, ROLL, SLIME, Seer"],
    ["GRPO (无Critic)", "组内相对优势, 消除Critic", "verl, ROLL, SLIME, MILES, Seer"],
    ["DAPO", "解耦Clip+动态采样, 防熵坍塌", "verl, ROLL, SLIME"],
    ["CISPO", "裁剪IS权重, 全Token梯度", "Forge"],
    ["SAPO", "共享Rollout文本(非梯度)", "rl-swarm"],
    ["Chunked MDP", "Chunk级IS, 匹配Agent交互", "ROLL"],
    ["DPO/SimPO", "无RL循环, 偏好分类", "verl (辅助)"],
]
cw = [Inches(2.0), Inches(4.5), Inches(4.5)]
_make_table(s, 8, 3, data, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.2),
            col_widths=cw, font_size=11)

# References
refs = [
    "1. verl/HybridFlow: arXiv 2409.19256, EuroSys 2025",
    "2. ROLL: arXiv 2506.06122 / ROLLMUX: 2512.11306 / ROLLART: 2512.22560",
    "3. SLIME/APRIL: arXiv 2509.18521",
    "4. MILES/R3: arXiv 2510.11370",
    "5. MiniMax Forge: HuggingFace Blog",
    "6. Seer: arXiv 2511.14617  |  7. rl-swarm/SAPO: arXiv 2509.08721",
    "8. ThunderAgent: arXiv 2602.13692 (Together AI)",
    "9. The Landscape of Agentic RL for LLMs: arXiv 2509.02547, TMLR 2026",
]
_add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
          "参考文献", size=16, color=ACCENT, bold=True)
tf = _add_text(s, Inches(0.5), Inches(5.0), Inches(12), Inches(2.2),
               refs[0], size=10, color=GRAY)
for ref in refs[1:]:
    _add_para(tf, ref, size=10, color=GRAY, space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
out = "/Users/xuelei/go/src/rl-research/agentic-rl-infra.pptx"
prs.save(out)
print(f"✅ PPT saved to {out} ({len(prs.slides)} slides)")
